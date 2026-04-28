import io
import json
import os

from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from openai import AsyncOpenAI
from pydantic import BaseModel

router = APIRouter()


def get_client() -> AsyncOpenAI:
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        raise HTTPException(
            status_code=503,
            detail="OpenAI API key not configured. Set the OPENAI_API_KEY environment variable.",
        )
    return AsyncOpenAI(api_key=api_key)


class WordDocRequest(BaseModel):
    title: str
    content_prompt: str
    style: str = "professional"


class ExcelRequest(BaseModel):
    title: str
    data_description: str
    num_rows: int = 10


class PptRequest(BaseModel):
    title: str
    topic: str
    num_slides: int = 8
    style: str = "professional"


@router.post("/generate-word")
async def generate_word(request: WordDocRequest):
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt, RGBColor

    client = get_client()

    prompt = f"""Generate content for a {request.style} Word document.

Title: {request.title}
Topic/Instructions: {request.content_prompt}

Return a JSON object with this structure:
{{
    "title": "Document Title",
    "sections": [
        {{
            "heading": "Section Heading",
            "content": "Section content paragraph(s). Use \\n for new paragraphs."
        }}
    ]
}}

Create detailed, professional content with at least 4-6 sections.
Return ONLY the JSON, no other text."""

    response = await client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "system",
                "content": "You are a document writer. Generate structured content for professional documents. Return only valid JSON.",
            },
            {"role": "user", "content": prompt},
        ],
        max_tokens=4000,
        temperature=0.7,
        response_format={"type": "json_object"},
    )

    try:
        doc_data = json.loads(response.choices[0].message.content)
    except json.JSONDecodeError:
        raise HTTPException(status_code=500, detail="Failed to generate document content")

    doc = Document()

    style = doc.styles["Title"]
    style.font.color.rgb = RGBColor(0x1A, 0x56, 0xDB)
    style.font.size = Pt(28)

    title_para = doc.add_heading(doc_data.get("title", request.title), level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph("")

    for section in doc_data.get("sections", []):
        heading = doc.add_heading(section.get("heading", ""), level=1)
        heading_run = heading.runs[0] if heading.runs else None
        if heading_run:
            heading_run.font.color.rgb = RGBColor(0x2D, 0x3A, 0x4A)

        content = section.get("content", "")
        for paragraph_text in content.split("\n"):
            if paragraph_text.strip():
                para = doc.add_paragraph(paragraph_text.strip())
                para.paragraph_format.space_after = Pt(6)
                for run in para.runs:
                    run.font.size = Pt(11)
                    run.font.name = "Calibri"

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    filename = request.title.replace(" ", "_").lower() + ".docx"

    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/generate-excel")
async def generate_excel(request: ExcelRequest):
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    client = get_client()

    prompt = f"""Generate spreadsheet data for: {request.data_description}

Title: {request.title}
Number of rows: {request.num_rows}

Return a JSON object with this structure:
{{
    "title": "Spreadsheet Title",
    "headers": ["Column1", "Column2", "Column3"],
    "rows": [
        ["value1", "value2", "value3"],
        ["value1", "value2", "value3"]
    ],
    "chart_title": "Chart Title (if applicable)"
}}

Make the data realistic and relevant. Include numeric columns where possible for charting.
Return ONLY the JSON, no other text."""

    response = await client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "system",
                "content": "You are a data analyst. Generate realistic spreadsheet data. Return only valid JSON. Ensure numeric values are actual numbers, not strings.",
            },
            {"role": "user", "content": prompt},
        ],
        max_tokens=4000,
        temperature=0.7,
        response_format={"type": "json_object"},
    )

    try:
        data = json.loads(response.choices[0].message.content)
    except json.JSONDecodeError:
        raise HTTPException(status_code=500, detail="Failed to generate spreadsheet data")

    wb = Workbook()
    ws = wb.active
    ws.title = data.get("title", request.title)[:31]

    header_font = Font(bold=True, color="FFFFFF", size=12)
    header_fill = PatternFill(start_color="1A56DB", end_color="1A56DB", fill_type="solid")
    header_alignment = Alignment(horizontal="center", vertical="center")
    thin_border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin"),
    )

    headers = data.get("headers", [])
    for col_idx, header in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col_idx, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = header_alignment
        cell.border = thin_border
        ws.column_dimensions[cell.column_letter].width = max(15, len(str(header)) + 5)

    rows = data.get("rows", [])
    for row_idx, row_data in enumerate(rows, 2):
        for col_idx, value in enumerate(row_data, 1):
            if isinstance(value, str):
                try:
                    value = float(value)
                    if value == int(value):
                        value = int(value)
                except (ValueError, TypeError):
                    pass
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.border = thin_border
            cell.alignment = Alignment(horizontal="center")

    numeric_col = None
    for col_idx, header in enumerate(headers, 1):
        if col_idx > 1:
            sample_value = ws.cell(row=2, column=col_idx).value
            if isinstance(sample_value, (int, float)):
                numeric_col = col_idx
                break

    if numeric_col and len(rows) > 1:
        chart = BarChart()
        chart.title = data.get("chart_title", f"{request.title} Chart")
        chart.style = 10
        chart.y_axis.title = headers[numeric_col - 1] if numeric_col <= len(headers) else "Values"
        chart.x_axis.title = headers[0] if headers else "Category"

        chart_data = Reference(ws, min_col=numeric_col, min_row=1, max_row=len(rows) + 1)
        cats = Reference(ws, min_col=1, min_row=2, max_row=len(rows) + 1)
        chart.add_data(chart_data, titles_from_data=True)
        chart.set_categories(cats)
        chart.shape = 4
        ws.add_chart(chart, "A" + str(len(rows) + 4))

    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)

    filename = request.title.replace(" ", "_").lower() + ".xlsx"

    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/generate-ppt")
async def generate_ppt(request: PptRequest):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    client = get_client()

    prompt = f"""Create a presentation on: {request.topic}

Title: {request.title}
Number of slides: {request.num_slides}
Style: {request.style}

Return a JSON object with this structure:
{{
    "title": "{request.title}",
    "subtitle": "Subtitle text",
    "slides": [
        {{
            "title": "Slide Title",
            "bullet_points": ["Point 1", "Point 2", "Point 3"],
            "notes": "Speaker notes for this slide"
        }}
    ]
}}

Make the content informative and well-structured. Each slide should have 3-5 bullet points.
Return ONLY the JSON, no other text."""

    response = await client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "system",
                "content": "You are a presentation expert. Create engaging, well-structured presentations. Return only valid JSON.",
            },
            {"role": "user", "content": prompt},
        ],
        max_tokens=4000,
        temperature=0.7,
        response_format={"type": "json_object"},
    )

    try:
        ppt_data = json.loads(response.choices[0].message.content)
    except json.JSONDecodeError:
        raise HTTPException(
            status_code=500, detail="Failed to generate presentation content"
        )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = ppt_data.get("title", request.title)
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1A, 0x56, 0xDB)

    if slide.placeholders[1]:
        subtitle = slide.placeholders[1]
        subtitle.text = ppt_data.get("subtitle", "")
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(20)
                run.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)

    # Content slides
    for slide_data in ppt_data.get("slides", []):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "")
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(32)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x1A, 0x56, 0xDB)

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        tf.word_wrap = True

        bullet_points = slide_data.get("bullet_points", [])
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.space_after = Pt(12)
            p.level = 0
            for run in p.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x37, 0x41, 0x51)

        notes = slide_data.get("notes", "")
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    # Thank you slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Thank You!"
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1A, 0x56, 0xDB)

    if slide.placeholders[1]:
        subtitle = slide.placeholders[1]
        subtitle.text = "Questions & Discussion"
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    filename = request.title.replace(" ", "_").lower() + ".pptx"

    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
